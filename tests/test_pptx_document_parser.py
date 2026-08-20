from __future__ import annotations

from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import pytest
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from filelore.documents import PptxDocumentParser, TextBlockType
from filelore.metadata import DocumentFormat


def create_text_pptx(path: Path) -> None:
    presentation = Presentation()
    core = presentation.core_properties
    core.title = "Türkçe PPTX Belgesi"
    core.author = "Örnek Yazar"
    core.subject = "Yapısal slayt metni çıkarma"
    core.keywords = "Türkçe, sunum, İstanbul"
    core.language = "tr"
    core.last_modified_by = "FileLore test suite"
    core.created = datetime(2024, 2, 3, 10, 30)
    core.modified = datetime(2024, 2, 4, 11, 45)

    first_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    first_slide.shapes.title.text = "Ana Slayt"
    text_frame = first_slide.placeholders[1].text_frame
    text_frame.clear()
    text_frame.paragraphs[0].text = (
        "Türkçe, 日本語 ve العربية aynı slaytta kalır."
    )
    first_item = text_frame.add_paragraph()
    first_item.text = "Birinci madde"
    _set_bullet(first_item)
    second_item = text_frame.add_paragraph()
    second_item.text = "İkinci madde"
    second_item.level = 1
    _set_bullet(second_item)

    table_shape = first_slide.shapes.add_table(
        2,
        2,
        Inches(0.7),
        Inches(5.4),
        Inches(6.0),
        Inches(1.0),
    )
    table = table_shape.table
    table.cell(0, 0).text = "Alan"
    table.cell(0, 1).text = "Değer"
    table.cell(1, 0).text = "Dil"
    table.cell(1, 1).text = "Türkçe"
    first_slide.notes_slide.notes_text_frame.text = (
        "Konuşmacı notu: çok dilli içerik korunur."
    )

    second_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    second_slide.shapes.title.text = "İkinci Slayt"
    text_box = second_slide.shapes.add_textbox(
        Inches(1.0),
        Inches(2.0),
        Inches(8.0),
        Inches(1.0),
    )
    text_box.text = "Bağımsız metin kutusu içeriği."
    presentation.save(path)


def _set_bullet(paragraph: Any) -> None:
    paragraph_properties = paragraph._p.get_or_add_pPr()
    bullet = etree.Element(qn("a:buChar"))
    bullet.set("char", "•")
    paragraph_properties.append(bullet)


def test_pptx_parser_extracts_slides_tables_notes_metadata_and_unicode(
    tmp_path: Path,
) -> None:
    path = tmp_path / "sunum.PPTX"
    create_text_pptx(path)

    document = PptxDocumentParser().parse(path)

    assert document.metadata.document_format is DocumentFormat.PPTX
    assert document.metadata.extension == ".pptx"
    assert document.metadata.mime_type.endswith("presentationml.presentation")
    assert document.metadata.title == "Türkçe PPTX Belgesi"
    assert document.metadata.authors == ("Örnek Yazar",)
    assert document.metadata.language == "tr"
    assert document.metadata.created_at == datetime(
        2024, 2, 3, 10, 30, tzinfo=timezone.utc
    )
    assert document.metadata.content_modified_at == datetime(
        2024, 2, 4, 11, 45, tzinfo=timezone.utc
    )
    assert document.metadata.slide_count == 2
    assert document.metadata.page_count is None
    assert document.metadata.properties["subject"] == (
        "Yapısal slayt metni çıkarma"
    )
    assert document.metadata.properties["keywords"] == (
        "Türkçe",
        "sunum",
        "İstanbul",
    )
    assert document.metadata.properties["last_modified_by"] == (
        "FileLore test suite"
    )
    assert document.metadata.properties["slide_width_emu"] > 0
    assert [block.block_type for block in document.blocks] == [
        TextBlockType.SLIDE_TITLE,
        TextBlockType.PARAGRAPH,
        TextBlockType.LIST_ITEM,
        TextBlockType.LIST_ITEM,
        TextBlockType.TABLE,
        TextBlockType.SPEAKER_NOTE,
        TextBlockType.SLIDE_TITLE,
        TextBlockType.PARAGRAPH,
    ]
    assert document.blocks[1].text == (
        "Türkçe, 日本語 ve العربية aynı slaytta kalır."
    )
    assert document.blocks[1].location.slide_number == 1
    assert document.blocks[1].location.section_path == ("Ana Slayt",)
    assert document.blocks[2].attributes["paragraph_level"] == 0
    assert document.blocks[3].attributes["paragraph_level"] == 1
    assert document.blocks[4].text == "Alan | Değer\nDil | Türkçe"
    assert document.blocks[4].attributes["rows"] == 2
    assert document.blocks[5].text.startswith("Konuşmacı notu")
    assert document.blocks[6].location.slide_number == 2
    assert document.blocks[7].text == "Bağımsız metin kutusu içeriği."


def test_pptx_parser_uses_the_first_slide_title_as_fallback(
    tmp_path: Path,
) -> None:
    path = tmp_path / "fallback.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Sunum Başlığı"
    presentation.core_properties.title = ""
    presentation.save(path)

    document = PptxDocumentParser().parse(path)

    assert document.metadata.title == "Sunum Başlığı"


def test_pptx_parser_rejects_a_presentation_without_slides(
    tmp_path: Path,
) -> None:
    path = tmp_path / "empty.pptx"
    Presentation().save(path)

    with pytest.raises(ValueError, match="contains no slides"):
        PptxDocumentParser().parse(path)


def test_pptx_parser_enforces_archive_limits(tmp_path: Path) -> None:
    path = tmp_path / "limited.pptx"
    create_text_pptx(path)

    with pytest.raises(ValueError, match="uncompressed size limit"):
        PptxDocumentParser(max_uncompressed_bytes=1).parse(path)
    with pytest.raises(ValueError, match="more than 1 members"):
        PptxDocumentParser(max_archive_members=1).parse(path)


def test_pptx_parser_rejects_corrupt_unsupported_and_missing_files(
    tmp_path: Path,
) -> None:
    corrupt_path = tmp_path / "corrupt.pptx"
    corrupt_path.write_text("not a PPTX", encoding="utf-8")
    parser = PptxDocumentParser()

    with pytest.raises(ValueError, match="Could not parse PPTX"):
        parser.parse(corrupt_path)
    with pytest.raises(ValueError, match="Unsupported document extension"):
        parser.parse(tmp_path / "notes.txt")
    with pytest.raises(FileNotFoundError):
        parser.parse(tmp_path / "missing.pptx")


@pytest.mark.parametrize(
    ("uncompressed_limit", "member_limit"),
    ((0, 1), (1, 0)),
)
def test_pptx_parser_rejects_invalid_archive_limits(
    uncompressed_limit: int,
    member_limit: int,
) -> None:
    with pytest.raises(ValueError, match="limit must be positive"):
        PptxDocumentParser(
            max_uncompressed_bytes=uncompressed_limit,
            max_archive_members=member_limit,
        )
