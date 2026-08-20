"""Slide-aware PPTX parser backed by python-pptx."""

from __future__ import annotations

from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable, Iterator
from zipfile import BadZipFile, ZipFile

from lxml.etree import XMLSyntaxError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import InvalidXmlError, PackageNotFoundError, PythonPptxError
from pptx.oxml.ns import qn
from pptx.presentation import Presentation as PptxPresentation
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from filelore.documents.models import (
    ParsedDocument,
    SourceLocation,
    TextBlock,
    TextBlockType,
)
from filelore.documents.parsers.base import DocumentParser
from filelore.metadata import DocumentFormat, DocumentMetadata


class PptxDocumentParser(DocumentParser):
    """Extract slide text, tables, notes, and normalized core properties."""

    supported_extensions = frozenset({".pptx"})
    default_max_uncompressed_bytes = 256 * 1024 * 1024
    default_max_archive_members = 10_000

    def __init__(
        self,
        *,
        max_uncompressed_bytes: int = default_max_uncompressed_bytes,
        max_archive_members: int = default_max_archive_members,
    ) -> None:
        if max_uncompressed_bytes <= 0:
            raise ValueError("PPTX uncompressed size limit must be positive")
        if max_archive_members <= 0:
            raise ValueError("PPTX archive member limit must be positive")
        self.max_uncompressed_bytes = max_uncompressed_bytes
        self.max_archive_members = max_archive_members

    def parse(self, path: str | Path) -> ParsedDocument:
        presentation_path = self.prepare_path(path)
        try:
            _validate_archive(
                presentation_path,
                max_uncompressed_bytes=self.max_uncompressed_bytes,
                max_archive_members=self.max_archive_members,
            )
        except BadZipFile as error:
            raise ValueError(
                f"Could not parse PPTX document: {presentation_path}"
            ) from error

        try:
            presentation = Presentation(presentation_path)
        except (
            BadZipFile,
            InvalidXmlError,
            PackageNotFoundError,
            PythonPptxError,
            XMLSyntaxError,
            KeyError,
            TypeError,
            ValueError,
        ) as error:
            raise ValueError(
                f"Could not parse PPTX document: {presentation_path}"
            ) from error

        if not presentation.slides:
            raise ValueError(
                f"PPTX document contains no slides: {presentation_path}"
            )

        try:
            blocks = _extract_blocks(presentation)
            metadata = _document_metadata(
                presentation,
                presentation_path,
                blocks,
            )
            return ParsedDocument(metadata=metadata, blocks=blocks)
        except (
            KeyError,
            TypeError,
            ValueError,
        ) as error:
            raise ValueError(
                f"Could not parse PPTX document: {presentation_path}"
            ) from error


def _validate_archive(
    path: Path,
    *,
    max_uncompressed_bytes: int,
    max_archive_members: int,
) -> None:
    with ZipFile(path) as archive:
        members = archive.infolist()
        if len(members) > max_archive_members:
            raise ValueError(
                f"PPTX archive contains more than {max_archive_members} members"
            )
        uncompressed_size = sum(member.file_size for member in members)
        if uncompressed_size > max_uncompressed_bytes:
            raise ValueError(
                "PPTX archive exceeds the "
                f"{max_uncompressed_bytes}-byte uncompressed size limit"
            )
        if any(member.flag_bits & 0x1 for member in members):
            raise ValueError("Encrypted PPTX archives are not supported")

        names = {member.filename for member in members}
        required = {"[Content_Types].xml", "ppt/presentation.xml"}
        if not required.issubset(names):
            raise ValueError("PPTX archive is missing required presentation parts")


def _extract_blocks(presentation: PptxPresentation) -> tuple[TextBlock, ...]:
    blocks: list[TextBlock] = []
    for slide_index, slide in enumerate(presentation.slides):
        slide_number = slide_index + 1
        title_shape = slide.shapes.title
        title_element = title_shape.element if title_shape is not None else None
        slide_title = (
            _normalize_text(title_shape.text)
            if title_shape is not None and title_shape.has_text_frame
            else ""
        )
        location = SourceLocation(
            slide_number=slide_number,
            section_path=(slide_title,) if slide_title else (),
        )

        for shape in _ordered_shapes(slide.shapes):
            if shape.has_table:
                _append_table_block(blocks, shape, location)
            elif shape.has_text_frame:
                _append_text_blocks(
                    blocks,
                    shape,
                    location,
                    is_title=shape.element is title_element,
                )

        _append_speaker_notes(blocks, slide, location)
    return tuple(blocks)


def _ordered_shapes(shapes: Iterable[BaseShape]) -> Iterator[BaseShape]:
    indexed = list(enumerate(shapes))
    indexed.sort(
        key=lambda item: (
            _coordinate(item[1], "top"),
            _coordinate(item[1], "left"),
            item[0],
        )
    )
    for _, shape in indexed:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _ordered_shapes(shape.shapes)
        else:
            yield shape


def _append_text_blocks(
    blocks: list[TextBlock],
    shape: BaseShape,
    location: SourceLocation,
    *,
    is_title: bool,
) -> None:
    attributes = _shape_attributes(shape)
    if is_title:
        text = _normalize_text(shape.text)
        if text:
            blocks.append(
                TextBlock(
                    index=len(blocks),
                    block_type=TextBlockType.SLIDE_TITLE,
                    text=text,
                    location=location,
                    attributes=attributes,
                )
            )
        return

    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
        text = _normalize_text(paragraph.text)
        if not text:
            continue
        paragraph_attributes = {
            **attributes,
            "paragraph_index": paragraph_index,
            "paragraph_level": paragraph.level,
        }
        blocks.append(
            TextBlock(
                index=len(blocks),
                block_type=(
                    TextBlockType.LIST_ITEM
                    if _is_list_paragraph(paragraph)
                    else TextBlockType.PARAGRAPH
                ),
                text=text,
                location=location,
                attributes=paragraph_attributes,
            )
        )


def _append_table_block(
    blocks: list[TextBlock],
    shape: BaseShape,
    location: SourceLocation,
) -> None:
    table = shape.table
    rows: list[str] = []
    for row in table.rows:
        cells = [_single_line_text(cell.text) for cell in row.cells]
        if any(cells):
            rows.append(" | ".join(cells))
    text = "\n".join(rows)
    if not text:
        return
    blocks.append(
        TextBlock(
            index=len(blocks),
            block_type=TextBlockType.TABLE,
            text=text,
            location=location,
            attributes={
                **_shape_attributes(shape),
                "rows": len(table.rows),
                "columns": len(table.columns),
            },
        )
    )


def _append_speaker_notes(
    blocks: list[TextBlock],
    slide: Slide,
    location: SourceLocation,
) -> None:
    if not slide.has_notes_slide:
        return
    text_frame = slide.notes_slide.notes_text_frame
    if text_frame is None:
        return
    for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
        text = _normalize_text(paragraph.text)
        if not text:
            continue
        blocks.append(
            TextBlock(
                index=len(blocks),
                block_type=TextBlockType.SPEAKER_NOTE,
                text=text,
                location=location,
                attributes={"paragraph_index": paragraph_index},
            )
        )


def _shape_attributes(shape: BaseShape) -> dict[str, Any]:
    attributes: dict[str, Any] = {}
    shape_name = _first_text(shape.name)
    if shape_name:
        attributes["shape_name"] = shape_name
    if shape.is_placeholder:
        placeholder_type = shape.placeholder_format.type
        attributes["placeholder_type"] = placeholder_type.name.casefold()
    return attributes


def _is_list_paragraph(paragraph: Any) -> bool:
    if paragraph.level > 0:
        return True
    paragraph_properties = paragraph._p.pPr
    if paragraph_properties is None:
        return False
    if paragraph_properties.find(qn("a:buNone")) is not None:
        return False
    return any(
        paragraph_properties.find(qn(tag)) is not None
        for tag in ("a:buAutoNum", "a:buBlip", "a:buChar")
    )


def _document_metadata(
    presentation: PptxPresentation,
    path: Path,
    blocks: tuple[TextBlock, ...],
) -> DocumentMetadata:
    core = presentation.core_properties
    author = _first_text(core.author)
    stat = path.stat()
    return DocumentMetadata(
        path=path,
        extension=path.suffix.casefold(),
        mime_type=(
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation"
        ),
        size_bytes=stat.st_size,
        modified_at=datetime.fromtimestamp(stat.st_mtime).astimezone(),
        document_format=DocumentFormat.PPTX,
        title=_first_text(core.title, _first_slide_title(blocks)),
        authors=(author,) if author else (),
        created_at=_core_datetime(core.created),
        content_modified_at=_core_datetime(core.modified),
        slide_count=len(presentation.slides),
        language=_first_text(core.language),
        properties=_document_properties(presentation),
    )


def _document_properties(presentation: PptxPresentation) -> dict[str, Any]:
    core = presentation.core_properties
    properties: dict[str, Any] = {
        "slide_width_emu": int(presentation.slide_width),
        "slide_height_emu": int(presentation.slide_height),
    }
    for key in (
        "category",
        "comments",
        "content_status",
        "identifier",
        "last_modified_by",
        "subject",
        "version",
    ):
        value = _first_text(getattr(core, key, None))
        if value:
            properties[key] = value
    keywords = _keywords(core.keywords)
    if keywords:
        properties["keywords"] = keywords
    if isinstance(core.revision, int):
        properties["revision"] = core.revision
    last_printed = _core_datetime(core.last_printed)
    if last_printed is not None:
        properties["last_printed"] = last_printed
    return properties


def _first_slide_title(blocks: tuple[TextBlock, ...]) -> str | None:
    for block in blocks:
        if block.block_type is TextBlockType.SLIDE_TITLE:
            return block.text
    return None


def _core_datetime(value: Any) -> datetime | None:
    if not isinstance(value, datetime):
        return None
    if value.tzinfo is None:
        return value.replace(tzinfo=timezone.utc)
    return value


def _coordinate(shape: BaseShape, name: str) -> int:
    value = getattr(shape, name, None)
    return int(value) if isinstance(value, int) else 0


def _normalize_text(value: str) -> str:
    normalized = (
        value.replace("\x00", "")
        .replace("\v", "\n")
        .replace("\r\n", "\n")
        .replace("\r", "\n")
    )
    lines = (_single_line_text(line) for line in normalized.splitlines())
    return "\n".join(line for line in lines if line)


def _single_line_text(value: str) -> str:
    return " ".join(value.replace("\x00", "").split())


def _first_text(*values: Any) -> str | None:
    for value in values:
        if value is None:
            continue
        text = _single_line_text(str(value))
        if text:
            return text
    return None


def _keywords(value: Any) -> tuple[str, ...]:
    if value is None:
        return ()
    prepared: list[str] = []
    seen: set[str] = set()
    for candidate in str(value).replace(";", ",").split(","):
        keyword = _first_text(candidate)
        if keyword is None or keyword.casefold() in seen:
            continue
        seen.add(keyword.casefold())
        prepared.append(keyword)
    return tuple(prepared)
